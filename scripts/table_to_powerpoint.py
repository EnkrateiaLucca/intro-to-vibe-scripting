#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = ["pandas", "matplotlib", "python-pptx"]
# ///

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Path to CSV file
csv_path = "/Users/greatmaster/Desktop/projects/content_creation/intro-to-vibe-scripting/assets/financial-data.csv"

# Load data into a DataFrame
df = pd.read_csv(csv_path)
df["Date"] = pd.to_datetime(df["Date"])

# Use the first row for analysis
row = df.iloc[0]
insights = {
    "Date": row["Date"].date(),
    "Opening Price": row["Open"],
    "Closing Price": row["Close"],
    "Daily High": row["High"],
    "Daily Low": row["Low"],
    "Volume Traded": int(row["Volume"]),
    "Daily Range": row["High"] - row["Low"],
    "Price Change": row["Close"] - row["Open"]
}

# Create a bar chart
plt.figure(figsize=(6, 4))
plt.bar(["Open", "High", "Low", "Close"],
        [insights["Opening Price"], insights["Daily High"], insights["Daily Low"], insights["Closing Price"]],
        color=["blue", "green", "red", "purple"])
plt.title("Stock Prices on {}".format(insights["Date"]))
plt.ylabel("Price ($)")

# Save chart to BytesIO
img_stream = BytesIO()
plt.savefig(img_stream, format='png')
img_stream.seek(0)

# Create PowerPoint
prs = Presentation()
slide_layout = prs.slide_layouts[5]

# Slide 1 - Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title, subtitle = slide.shapes.title, slide.placeholders[1]
title.text = "Daily Stock Report"
subtitle.text = f"Insights for {insights['Date']}"

# Slide 2 - Insights
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
shapes.title.text = "Key Insights"

insight_text = "\n".join([
    f"Opening Price: ${insights['Opening Price']:.2f}",
    f"Closing Price: ${insights['Closing Price']:.2f}",
    f"Daily High: ${insights['Daily High']:.2f}",
    f"Daily Low: ${insights['Daily Low']:.2f}",
    f"Volume Traded: {insights['Volume Traded']:,}",
    f"Daily Range: ${insights['Daily Range']:.2f}",
    f"Price Change: ${insights['Price Change']:.2f}"
])
text_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = insight_text

# Slide 3 - Chart
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
shapes.title.text = "Stock Price Chart"

pic = slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(6))

# Save the presentation
prs.save("stock_report.pptx")

print("Presentation saved as 'stock_report.pptx'")
    